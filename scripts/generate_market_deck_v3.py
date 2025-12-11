"""
Generate Agent9 Market Penetration PowerPoint Deck v3 (Widescreen + Graphics)

Usage:
    python scripts/generate_market_deck_v3.py

Output:
    docs/strategy/Agent9_Market_Penetration_Deck_v3.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path
import math

# Brand colors
NAVY = RGBColor(0x1a, 0x36, 0x5d)  # Dark navy blue
BLUE = RGBColor(0x2c, 0x52, 0x82)  # Medium blue
LIGHT_BLUE = RGBColor(0xeb, 0xf8, 0xff) # Very light blue background
WHITE = RGBColor(0xff, 0xff, 0xff)
ACCENT_GREEN = RGBColor(0x2e, 0x7d, 0x32)  # Green for highlights
ACCENT_ORANGE = RGBColor(0xdd, 0x6b, 0x20) # Orange for accents
LIGHT_GRAY = RGBColor(0xf5, 0xf5, 0xf5)
DARK_GRAY = RGBColor(0x42, 0x42, 0x42)

# Dimensions (16:9 Widescreen)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

def set_slide_background(slide, color):
    """Set solid background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(prs, title, subtitle, tagline=None):
    """Add a modern title slide with a graphic element."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, NAVY)
    
    # Graphic Element (Side Bar)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(0), Inches(0), Inches(0.5), SLIDE_HEIGHT
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_GREEN
    shape.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(11), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(1.5), Inches(3.6), Inches(11), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(32)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    
    # Tagline
    if tagline:
        tag_box = slide.shapes.add_textbox(Inches(1.5), Inches(5.0), Inches(11), Inches(1))
        tf = tag_box.text_frame
        p = tf.paragraphs[0]
        p.text = tagline
        p.font.size = Pt(24)
        p.font.italic = True
        p.font.color.rgb = LIGHT_GRAY
        p.alignment = PP_ALIGN.LEFT
    
    return slide

def add_header(slide, title):
    """Add consistent header to a slide."""
    # Title bar
    title_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.2))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = NAVY
    title_shape.line.fill.background()
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

def add_content_slide(prs, title, bullets, notes=None):
    """Add a standard content slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, WHITE)
    add_header(slide, title)
    
    bullet_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(5))
    tf = bullet_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(24)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(14)

    if notes: slide.notes_slide.notes_text_frame.text = notes
    return slide

def add_pipeline_visual_slide(prs, title, notes=None):
    """Add the Intelligence Pipeline visual (Chevrons)."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, WHITE)
    add_header(slide, title)

    # Chevron 1: Situation Awareness
    shape1 = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(1), Inches(2.5), Inches(3.5), Inches(2))
    shape1.fill.solid()
    shape1.fill.fore_color.rgb = BLUE
    shape1.line.color.rgb = WHITE
    shape1.text_frame.text = "1. SITUATION\nAWARENESS"
    shape1.text_frame.paragraphs[0].font.size = Pt(20)
    shape1.text_frame.paragraphs[0].font.bold = True
    shape1.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Chevron 2: Deep Analysis
    shape2 = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(4.8), Inches(2.5), Inches(3.5), Inches(2))
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = BLUE
    shape2.line.color.rgb = WHITE
    shape2.text_frame.text = "2. DEEP\nANALYSIS"
    shape2.text_frame.paragraphs[0].font.size = Pt(20)
    shape2.text_frame.paragraphs[0].font.bold = True
    shape2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Chevron 3: Solution Finding
    shape3 = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(8.6), Inches(2.5), Inches(3.5), Inches(2))
    shape3.fill.solid()
    shape3.fill.fore_color.rgb = ACCENT_GREEN
    shape3.line.color.rgb = WHITE
    shape3.text_frame.text = "3. SOLUTION\nFINDING"
    shape3.text_frame.paragraphs[0].font.size = Pt(20)
    shape3.text_frame.paragraphs[0].font.bold = True
    shape3.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Descriptions under chevrons
    def add_desc(x, text):
        box = slide.shapes.add_textbox(x, Inches(4.7), Inches(3.5), Inches(1.5))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(16)
        p.alignment = PP_ALIGN.CENTER
        p.font.color.rgb = DARK_GRAY

    add_desc(Inches(1), "• Continuous KPI monitoring\n• Anomaly detection\n• Impact prioritization")
    add_desc(Inches(4.8), "• Root cause investigation\n• Multi-dimensional drill-down\n• Evidence gathering")
    add_desc(Inches(8.6), "• Multi-agent debate\n• Branded expert perspectives\n• Consensus synthesis")

    # Bottom Tagline
    tag_box = slide.shapes.add_textbox(Inches(1), Inches(6.2), Inches(11.33), Inches(1))
    tf = tag_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Unlike ChatGPT, Agent9 continuously monitors your business and answers \"WHY?\" before you ask."
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER

    if notes: slide.notes_slide.notes_text_frame.text = notes
    return slide

def add_business_model_visual(prs, title, notes=None):
    """Add a Hub and Spoke visual for Business Model."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, WHITE)
    add_header(slide, title)

    cx, cy = SLIDE_WIDTH/2, Inches(4)
    
    # Center Hub (Agent9)
    center = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - Inches(1.5), cy - Inches(1.5), Inches(3), Inches(3))
    center.fill.solid()
    center.fill.fore_color.rgb = NAVY
    center.text_frame.text = "AGENT9\nPLATFORM"
    center.text_frame.paragraphs[0].font.bold = True
    center.text_frame.paragraphs[0].font.size = Pt(24)
    center.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Left Spoke (Partners)
    left_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), cy - Inches(1), Inches(3), Inches(2))
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = BLUE
    left_box.text_frame.text = "PARTNERS (Supply)\n• Consulting Firms\n• Domain Experts"
    
    # Right Spoke (Customers)
    right_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.33), cy - Inches(1), Inches(3), Inches(2))
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = BLUE
    right_box.text_frame.text = "CUSTOMERS (Demand)\n• Enterprises\n• PE Firms"

    # Arrows
    arrow_l = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.2), cy - Inches(0.25), Inches(1.5), Inches(0.5))
    arrow_l.fill.solid()
    arrow_l.fill.fore_color.rgb = LIGHT_GRAY
    
    arrow_r = slide.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, Inches(7.6), cy - Inches(0.25), Inches(1.5), Inches(0.5))
    arrow_r.fill.solid()
    arrow_r.fill.fore_color.rgb = LIGHT_GRAY

    # Revenue Stream boxes below
    rev_y = Inches(6)
    def add_rev_box(x, text):
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, rev_y, Inches(2.5), Inches(1))
        box.fill.solid()
        box.fill.fore_color.rgb = ACCENT_GREEN
        box.text_frame.text = text
        box.text_frame.paragraphs[0].font.size = Pt(14)
        box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_rev_box(Inches(1.5), "Platform Fees\n$80K-$300K/yr")
    add_rev_box(Inches(4.2), "Transaction Fees\n$100-$500/debate")
    add_rev_box(Inches(6.9), "Partner Share\n15-30%")
    add_rev_box(Inches(9.6), "Implementation\n$50K-$150K")

    if notes: slide.notes_slide.notes_text_frame.text = notes
    return slide

def add_matrix_slide(prs, title, notes=None):
    """Add a 2x2 Matrix visual."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, WHITE)
    add_header(slide, title)

    # Matrix Background
    left, top = Inches(3), Inches(2)
    width, height = Inches(7), Inches(5)
    
    # Axes Lines
    cx, cy = left + width/2, top + height/2
    
    # Vertical Axis Line
    vert = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, top, Inches(0.05), height)
    vert.fill.solid(); vert.fill.fore_color.rgb = DARK_GRAY
    
    # Horizontal Axis Line
    horz = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, cy, width, Inches(0.05))
    horz.fill.solid(); horz.fill.fore_color.rgb = DARK_GRAY

    # Labels
    def add_label(x, y, text, align):
        tb = slide.shapes.add_textbox(x, y, Inches(2), Inches(0.5))
        p = tb.text_frame.paragraphs[0]
        p.text = text
        p.font.bold = True
        p.alignment = align

    add_label(cx - Inches(1), top - Inches(0.5), "HIGH EXPERTISE", PP_ALIGN.CENTER)
    add_label(cx - Inches(1), top + height, "LOW EXPERTISE", PP_ALIGN.CENTER)
    add_label(left - Inches(2.1), cy - Inches(0.25), "LOW SCALE", PP_ALIGN.RIGHT)
    add_label(left + width + Inches(0.1), cy - Inches(0.25), "HIGH SCALE", PP_ALIGN.LEFT)

    # Quadrant Content
    def add_quad_content(x, y, title, subtitle, is_hero=False):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3), Inches(2))
        if is_hero:
            box.fill.solid(); box.fill.fore_color.rgb = NAVY
            box.text_frame.text = f"{title}\n\n{subtitle}"
            box.text_frame.paragraphs[0].font.color.rgb = WHITE
        else:
            box.fill.background()
            box.line.color.rgb = DARK_GRAY
            box.line.dash_style = 1 # dashed
            box.text_frame.text = f"{title}\n\n{subtitle}"
            box.text_frame.paragraphs[0].font.color.rgb = DARK_GRAY

    # Top Left: Traditional
    add_quad_content(left + Inches(0.2), top + Inches(0.2), "Traditional Consulting", "(McKinsey/BCG)\nHigh Expertise, Low Scale")
    # Bottom Right: Generic AI
    add_quad_content(left + width/2 + Inches(0.2), top + height/2 + Inches(0.2), "Generic AI", "(ChatGPT)\nHigh Scale, Low Expertise")
    # Top Right: Agent9 (Hero)
    add_quad_content(left + width/2 + Inches(0.2), top + Inches(0.2), "Agent9 CaaS", "Branded Agents at AI Scale\n★ WE ARE HERE", True)
    
    if notes: slide.notes_slide.notes_text_frame.text = notes
    return slide

def add_flywheel_slide(prs, title, points, notes=None):
    """Add a flywheel visual."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, WHITE)
    add_header(slide, title)

    cx, cy = SLIDE_WIDTH/2, Inches(4.5)
    radius = Inches(2.5)

    # Central Hub
    hub = slide.shapes.add_shape(MSO_SHAPE.DONUT, cx - Inches(1), cy - Inches(1), Inches(2), Inches(2))
    hub.fill.solid(); hub.fill.fore_color.rgb = NAVY
    hub.text_frame.text = "Moat"
    
    # Points around circle
    import math
    for i, (icon, text, desc) in enumerate(points):
        angle = 2 * math.pi * i / len(points) - math.pi/2 # Start at top
        x = cx + radius * math.cos(angle) - Inches(1.5)
        y = cy + radius * math.sin(angle) - Inches(0.75)
        
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3), Inches(1.5))
        box.fill.solid(); box.fill.fore_color.rgb = LIGHT_BLUE
        box.line.color.rgb = NAVY
        
        p = box.text_frame.paragraphs[0]
        p.text = f"{icon} {text}"
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = NAVY
        
        p2 = box.text_frame.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(12)

    if notes: slide.notes_slide.notes_text_frame.text = notes
    return slide

def add_table_slide(prs, title, headers, rows, notes=None):
    """Add a formatted table slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, WHITE)
    add_header(slide, title)
    
    rows_count = len(rows) + 1
    cols_count = len(headers)
    
    left = Inches(1)
    top = Inches(2)
    width = Inches(11.33)
    height = Inches(0.8 * rows_count)
    
    table = slide.shapes.add_table(rows_count, cols_count, left, top, width, height).table
    
    # Headers
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE
        cell.text_frame.paragraphs[0].font.color.rgb = WHITE
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
    # Rows
    for r, row_data in enumerate(rows):
        for c, item in enumerate(row_data):
            cell = table.cell(r+1, c)
            cell.text = str(item)
            cell.text_frame.paragraphs[0].font.size = Pt(16)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY

    if notes: slide.notes_slide.notes_text_frame.text = notes
    return slide

def add_two_column_slide(prs, title, left_title, left_bullets, right_title, right_bullets, notes=None):
    """Add a two-column slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, WHITE)
    add_header(slide, title)

    # Left Column
    lt = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5))
    lt.text_frame.text = left_title
    lt.text_frame.paragraphs[0].font.bold = True
    lt.text_frame.paragraphs[0].font.size = Pt(24)
    lt.text_frame.paragraphs[0].font.color.rgb = NAVY

    lb = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(5.5), Inches(4.5))
    l_tf = lb.text_frame
    l_tf.word_wrap = True
    for b in left_bullets:
        p = l_tf.add_paragraph()
        p.text = f"• {b}"
        p.font.size = Pt(20)
        p.space_after = Pt(10)

    # Right Column
    rt = slide.shapes.add_textbox(Inches(7), Inches(1.5), Inches(5.5), Inches(0.5))
    rt.text_frame.text = right_title
    rt.text_frame.paragraphs[0].font.bold = True
    rt.text_frame.paragraphs[0].font.size = Pt(24)
    rt.text_frame.paragraphs[0].font.color.rgb = NAVY

    rb = slide.shapes.add_textbox(Inches(7), Inches(2.2), Inches(5.5), Inches(4.5))
    r_tf = rb.text_frame
    r_tf.word_wrap = True
    for b in right_bullets:
        p = r_tf.add_paragraph()
        p.text = f"• {b}"
        p.font.size = Pt(20)
        p.space_after = Pt(10)

    if notes: slide.notes_slide.notes_text_frame.text = notes
    return slide

def create_presentation_v3():
    """Generate the V3 deck."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # 1. Title
    add_title_slide(prs, "Agent9", "The Agentic Consulting Marketplace", "Transforming $800B of consulting into on-demand, AI-powered expertise")

    # 2. Problem
    add_content_slide(prs, "Enterprise Decision-Making is Broken", [
        "⏱️  Consulting is slow — 12-24 weeks for strategic insights",
        "💰  Consulting is expensive — $500K-$2M minimum engagements",
        "👁️  Single perspective — One firm's bias, no debate",
        "📋  No audit trail — \"Trust us\" doesn't satisfy boards",
        "🚪  IP walks out — Knowledge leaves with consultants"
    ], "Let me paint the picture...")

    # 3. Opportunity
    add_table_slide(prs, "$800B Market, Zero Agentic Players", 
        ["Capability", "Traditional Consulting", "Generic AI (ChatGPT)", "Agent9 CaaS"],
        [
            ["Expertise", "✅", "❌", "✅"],
            ["Trust", "✅", "❌", "✅"],
            ["Speed", "❌", "✅", "✅"],
            ["Cost", "❌", "✅", "✅"],
            ["Audit Trail", "❌", "❌", "✅"]
        ], "Here is the gap in the market...")

    # 4. Solution
    add_content_slide(prs, "Agent9: Agentic Consulting Marketplace", [
        "Multiple branded AI agents debate your problem",
        "   → McKinsey Agent + BCG Agent + AWS Agent + ...",
        "Platform synthesizes best ideas into recommendation",
        "   → Multi-agent debate surfaces optimal solutions",
        "Complete audit trail for every decision",
        "   → Full transcript, evidence citations, confidence scores"
    ], "When a customer has a strategic question...")

    # 5. Intelligence Pipeline (NEW VISUAL)
    add_pipeline_visual_slide(prs, "The Intelligence Pipeline: How Agent9 Thinks", 
        "This is our core technical differentiation. Agent9 isn't a chatbot...")

    # 6. Pipeline Details (Situation Awareness)
    add_content_slide(prs, "Stage 1: Automated Situation Awareness", [
        "🔍  WHAT IT DOES",
        "   • Continuously monitors 100s of KPIs across your data",
        "   • Detects anomalies, trends, and threshold breaches",
        "   • Prioritizes by business impact and principal context",
        "",
        "⚡  WHY IT MATTERS",
        "   • Problems surface in hours, not quarterly reviews",
        "   • No analyst time spent on routine monitoring"
    ], "Situation Awareness is our always-on monitoring layer...")

    # 7. Pipeline Details (Deep Analysis)
    add_content_slide(prs, "Stage 2: Automated Deep Analysis", [
        "🔍  WHAT IT DOES",
        "   • Investigates root causes of detected situations",
        "   • Drills down across dimensions (region, product, channel)",
        "   • Builds evidence-based problem framing",
        "",
        "⚡  WHY IT MATTERS",
        "   • Answers \"why?\" before you have to ask",
        "   • Eliminates weeks of analyst investigation"
    ], "When Situation Awareness flags an issue, Deep Analysis investigates...")

    # 8. Pipeline Details (Solution Finding)
    add_content_slide(prs, "Stage 3: Multi-Agent Solution Finding", [
        "🔍  WHAT IT DOES",
        "   • Multiple expert agents debate solution options",
        "   • Each agent applies its methodology to the problem",
        "   • Platform synthesizes consensus recommendation",
        "",
        "⚡  WHY IT MATTERS",
        "   • Multiple perspectives, not single-firm bias",
        "   • Full debate transcript for audit/compliance"
    ], "This is where the CaaS marketplace comes alive...")

    # 9. Business Model (NEW VISUAL)
    add_business_model_visual(prs, "Three-Sided Marketplace", 
        "Our business model connects Partners (Supply) with Customers (Demand) via the Platform...")

    # 10. Target Market
    add_table_slide(prs, "Who We're Targeting", 
        ["Segment", "Market Size", "Avg. Deal", "Priority"],
        [
            ["PE Portfolio Ops", "500 firms", "$200K+", "⭐⭐⭐"],
            ["Mid-Market CFOs", "10,000+", "$80K", "⭐⭐⭐"],
            ["Corp Strategy Teams", "500 F500", "$300K+", "⭐⭐"],
            ["Data Platform Customers", "5,000+", "$100K", "⭐⭐"]
        ])

    # 11. GTM
    add_content_slide(prs, "Go-to-Market: Land, Expand, Partner", [
        "PHASE 1: LAND (Q1-Q2 2025)",
        "   • Direct sales to 10-20 early adopters",
        "   • Goal: 5-10 customers, $500K ARR",
        "",
        "PHASE 2: EXPAND (Q3-Q4 2025)",
        "   • Grow within accounts, add use cases",
        "   • Goal: $150K avg ACV, $2M ARR",
        "",
        "PHASE 3: PARTNER (2026)",
        "   • Launch branded agent marketplace",
        "   • Goal: 50% revenue from partner channel"
    ])

    # 12. Competitive Landscape (NEW VISUAL)
    add_matrix_slide(prs, "We're Creating a New Category", 
        "Agent9 occupies the upper right quadrant: Branded Expertise at AI Scale...")

    # 13. Defensibility (NEW VISUAL)
    add_flywheel_slide(prs, "Our Moat Deepens Over Time", [
        ("🤝", "Partner Network", "More expertise → more customers"),
        ("📊", "Data Network", "More customers → better patterns"),
        ("🔒", "Methodology", "IP encoded, hard to replicate"),
        ("📈", "Flywheel", "Network effects compound")
    ], "Our defensibility comes from network effects...")

    # 14. Financials
    add_table_slide(prs, "Path to $10M+ ARR", 
        ["Metric", "2025", "2026", "2027"],
        [
            ["Customers", "15", "50", "150"],
            ["Avg. ACV", "$100K", "$120K", "$130K"],
            ["Total Revenue", "$2M", "$7.5M", "$23M"]
        ])

    # 15. Team
    add_content_slide(prs, "The Team", [
        "[CEO NAME] - Consulting/Enterprise Software",
        "[CTO NAME] - AI/ML Platform Architecture",
        "[HEAD OF PARTNERSHIPS] - Consulting Alliances",
        "",
        "ADVISORS:",
        "• Former McKinsey Partner",
        "• Former CDO, Fortune 500"
    ])

    # 16. Ask
    add_two_column_slide(prs, "The Ask", 
        "Investors 💰", ["$X M Seed/Series A", "50% Engineering", "30% GTM", "20% Ops"],
        "Partners & Customers 🤝", ["2-3 Strategic Partners", "5-10 Early Adopters", "Pilot Programs"]
    )

    # 17. Summary
    add_content_slide(prs, "Summary: The Agentic Consulting Marketplace", [
        "💰 $800B market, underserved by technology",
        "⚡ 10-50x faster insights",
        "📉 70-90% cheaper than traditional firms",
        "🗣️ Multi-perspective AI debate",
        "✅ Full audit trail and compliance"
    ])

    return prs

if __name__ == "__main__":
    print("Generating Agent9 Market Penetration Deck v3 (Widescreen)...")
    prs = create_presentation_v3()
    
    script_dir = Path(__file__).parent
    output_path = script_dir.parent / "docs" / "strategy" / "Agent9_Market_Penetration_Deck_v3.pptx"
    
    prs.save(str(output_path))
    print(f"✅ Presentation saved to: {output_path}")
    print(f"   Total slides: {len(prs.slides)}")
