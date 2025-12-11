"""
Generate Agent9 Market Penetration PowerPoint Deck

Usage:
    python scripts/generate_market_deck.py

Output:
    docs/strategy/Agent9_Market_Penetration_Deck.pptx

Requirements:
    pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path


# Brand colors
NAVY = RGBColor(0x1a, 0x36, 0x5d)  # Dark navy blue
WHITE = RGBColor(0xff, 0xff, 0xff)
ACCENT_GREEN = RGBColor(0x2e, 0x7d, 0x32)  # Green for highlights
LIGHT_GRAY = RGBColor(0xf5, 0xf5, 0xf5)
DARK_GRAY = RGBColor(0x42, 0x42, 0x42)


def set_slide_background(slide, color):
    """Set solid background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_slide(prs, title, subtitle, tagline=None):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, NAVY)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(9), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Tagline
    if tagline:
        tag_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.6), Inches(9), Inches(0.6))
        tf = tag_box.text_frame
        p = tf.paragraphs[0]
        p.text = tagline
        p.font.size = Pt(18)
        p.font.italic = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    
    return slide


def add_content_slide(prs, title, bullets, notes=None):
    """Add a standard content slide with title and bullets."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, WHITE)
    
    # Title bar
    title_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = NAVY
    title_shape.line.fill.background()
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Bullets
    bullet_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    tf = bullet_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(24)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(12)
        p.level = 0
    
    # Speaker notes
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes
    
    return slide


def add_two_column_slide(prs, title, left_title, left_bullets, right_title, right_bullets, notes=None):
    """Add a two-column content slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, WHITE)
    
    # Title bar
    title_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = NAVY
    title_shape.line.fill.background()
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Left column title
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(4.2), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = NAVY
    
    # Left column bullets
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(4.2), Inches(4.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(left_bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(8)
    
    # Right column title
    right_title_box = slide.shapes.add_textbox(Inches(5.3), Inches(1.4), Inches(4.2), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = NAVY
    
    # Right column bullets
    right_box = slide.shapes.add_textbox(Inches(5.3), Inches(1.9), Inches(4.2), Inches(4.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(right_bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(8)
    
    # Speaker notes
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes
    
    return slide


def add_table_slide(prs, title, headers, rows, notes=None):
    """Add a slide with a table."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, WHITE)
    
    # Title bar
    title_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = NAVY
    title_shape.line.fill.background()
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Table
    num_cols = len(headers)
    num_rows = len(rows) + 1  # +1 for header
    
    table_width = Inches(9)
    table_height = Inches(0.5 * num_rows)
    left = Inches(0.5)
    top = Inches(1.6)
    
    table = slide.shapes.add_table(num_rows, num_cols, left, top, table_width, table_height).table
    
    # Set column widths
    col_width = table_width / num_cols
    for i in range(num_cols):
        table.columns[i].width = int(col_width)
    
    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    
    # Data rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(14)
            p.font.color.rgb = DARK_GRAY
            p.alignment = PP_ALIGN.CENTER
            # Alternate row colors
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
    
    # Speaker notes
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes
    
    return slide


def add_quote_slide(prs, quote, attribution):
    """Add a quote slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, NAVY)
    
    # Quote
    quote_box = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(8), Inches(2))
    tf = quote_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f'"{quote}"'
    p.font.size = Pt(32)
    p.font.italic = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Attribution
    attr_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(0.5))
    tf = attr_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"— {attribution}"
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    return slide


def add_section_slide(prs, title):
    """Add a section divider slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, NAVY)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    return slide


def add_key_points_slide(prs, title, points, notes=None):
    """Add a slide with key points (icon + text format)."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, WHITE)
    
    # Title bar
    title_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = NAVY
    title_shape.line.fill.background()
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Key points in two columns
    left_points = points[:3]
    right_points = points[3:6] if len(points) > 3 else []
    
    y_start = 1.6
    
    # Left column
    for i, (icon, text, desc) in enumerate(left_points):
        # Icon/emoji
        icon_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_start + i * 1.4), Inches(0.6), Inches(0.5))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.text = icon
        p.font.size = Pt(28)
        
        # Title
        text_box = slide.shapes.add_textbox(Inches(1.1), Inches(y_start + i * 1.4), Inches(3.5), Inches(0.4))
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = NAVY
        
        # Description
        desc_box = slide.shapes.add_textbox(Inches(1.1), Inches(y_start + 0.4 + i * 1.4), Inches(3.5), Inches(0.8))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
    
    # Right column
    for i, (icon, text, desc) in enumerate(right_points):
        # Icon/emoji
        icon_box = slide.shapes.add_textbox(Inches(5.3), Inches(y_start + i * 1.4), Inches(0.6), Inches(0.5))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.text = icon
        p.font.size = Pt(28)
        
        # Title
        text_box = slide.shapes.add_textbox(Inches(5.9), Inches(y_start + i * 1.4), Inches(3.5), Inches(0.4))
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = NAVY
        
        # Description
        desc_box = slide.shapes.add_textbox(Inches(5.9), Inches(y_start + 0.4 + i * 1.4), Inches(3.5), Inches(0.8))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
    
    # Speaker notes
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes
    
    return slide


def create_presentation():
    """Create the full presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # ========== SLIDE 1: Title ==========
    add_title_slide(
        prs,
        "Agent9",
        "The Agentic Consulting Marketplace",
        "Transforming $800B of consulting into on-demand, AI-powered expertise"
    )
    
    # ========== SLIDE 2: The Problem ==========
    add_content_slide(
        prs,
        "Enterprise Decision-Making is Broken",
        [
            "⏱️  Consulting is slow — 12-24 weeks for strategic insights",
            "💰  Consulting is expensive — $500K-$2M minimum engagements",
            "👁️  Single perspective — One firm's bias, no debate",
            "📋  No audit trail — \"Trust us\" doesn't satisfy boards",
            "🚪  IP walks out — Knowledge leaves with consultants"
        ],
        notes="Let me paint the picture of what enterprises face today. They're paying millions for insights that arrive months late. They get one firm's perspective with no way to verify the reasoning. And when the engagement ends, all that knowledge walks out the door. This is an $800 billion market ripe for disruption."
    )
    
    # ========== SLIDE 3: The Opportunity ==========
    add_table_slide(
        prs,
        "$800B Market, Zero Agentic Players",
        ["Capability", "Traditional Consulting", "AI Tools (ChatGPT)", "Agent9 CaaS"],
        [
            ["Expertise", "✅", "❌", "✅"],
            ["Trust", "✅", "❌", "✅"],
            ["Speed", "❌", "✅", "✅"],
            ["Cost", "❌", "✅", "✅"],
            ["Audit Trail", "❌", "❌", "✅"],
        ],
        notes="Here's the gap in the market. Traditional consulting has expertise but lacks speed and affordability. Generic AI tools are fast and cheap but lack real expertise and accountability. Agent9 sits in the middle—we deliver branded expertise at AI speed with full audit trails. No one else is doing this."
    )
    
    # ========== SLIDE 4: Our Solution ==========
    add_content_slide(
        prs,
        "Agent9: Agentic Consulting Marketplace",
        [
            "Multiple branded AI agents debate your problem",
            "   → McKinsey Agent + BCG Agent + AWS Agent + ...",
            "",
            "Platform synthesizes best ideas into recommendation",
            "   → Multi-agent debate surfaces optimal solutions",
            "",
            "Complete audit trail for every decision",
            "   → Full transcript, evidence citations, confidence scores"
        ],
        notes="Here's how Agent9 works. When a customer has a strategic question, they don't get one consultant's opinion—they get multiple branded AI agents, each trained on a firm's proprietary methodology, debating the problem. The platform synthesizes the best ideas and delivers a recommendation with a complete audit trail."
    )
    
    # ========== SLIDE 5: The Intelligence Pipeline ==========
    add_content_slide(
        prs,
        "The Intelligence Pipeline: How Agent9 Thinks",
        [
            "┌─────────────────────────────────────────────────────────────┐",
            "│  1. SITUATION        2. DEEP            3. SOLUTION         │",
            "│     AWARENESS   →      ANALYSIS    →      FINDING           │",
            "│                                                             │",
            "│  Continuous KPI      Root cause         Multi-agent         │",
            "│  monitoring          investigation      debate              │",
            "└─────────────────────────────────────────────────────────────┘",
            "",
            "Unlike ChatGPT: Agent9 doesn't wait for questions.",
            "It continuously monitors your business and surfaces issues",
            "BEFORE you know to ask."
        ],
        notes="This is our core technical differentiation. Agent9 isn't a chatbot waiting for questions—it's an always-on intelligence system. Situation Awareness continuously monitors KPIs. Deep Analysis investigates anomalies. Solution Finding debates options. This pipeline is what makes the consulting marketplace valuable."
    )
    
    # ========== SLIDE 6: Situation Awareness ==========
    add_content_slide(
        prs,
        "Stage 1: Automated Situation Awareness",
        [
            "🔍  WHAT IT DOES",
            "   • Continuously monitors 100s of KPIs across your data",
            "   • Detects anomalies, trends, and threshold breaches",
            "   • Prioritizes by business impact and principal context",
            "",
            "⚡  WHY IT MATTERS",
            "   • Problems surface in hours, not quarterly reviews",
            "   • No analyst time spent on routine monitoring",
            "   • Context-aware: knows what matters to each executive",
            "",
            "📊  EXAMPLE",
            "   \"Revenue in APAC dropped 12% vs. forecast—flagged for CFO\""
        ],
        notes="Situation Awareness is our always-on monitoring layer. It connects to your data warehouse and continuously scans KPIs. When something significant happens—a revenue drop, a margin squeeze, a customer churn spike—it surfaces immediately to the right executive. This replaces the army of analysts building dashboards that no one looks at."
    )
    
    # ========== SLIDE 7: Deep Analysis ==========
    add_content_slide(
        prs,
        "Stage 2: Automated Deep Analysis",
        [
            "🔍  WHAT IT DOES",
            "   • Investigates root causes of detected situations",
            "   • Drills down across dimensions (region, product, channel)",
            "   • Builds evidence-based problem framing",
            "",
            "⚡  WHY IT MATTERS",
            "   • Answers \"why?\" before you have to ask",
            "   • Eliminates weeks of analyst investigation",
            "   • Provides structured input for solution finding",
            "",
            "📊  EXAMPLE",
            "   \"APAC revenue drop driven by 3 factors: currency (40%),",
            "    delayed product launch (35%), competitor pricing (25%)\""
        ],
        notes="When Situation Awareness flags an issue, Deep Analysis automatically investigates. It doesn't just tell you revenue is down—it tells you WHY. It drills across every dimension in your data, quantifies contributing factors, and builds a structured problem frame. This is what consultants spend 4-8 weeks doing manually."
    )
    
    # ========== SLIDE 8: Solution Finding ==========
    add_content_slide(
        prs,
        "Stage 3: Multi-Agent Solution Finding",
        [
            "🔍  WHAT IT DOES",
            "   • Multiple expert agents debate solution options",
            "   • Each agent applies its methodology to the problem",
            "   • Platform synthesizes consensus recommendation",
            "",
            "⚡  WHY IT MATTERS",
            "   • Multiple perspectives, not single-firm bias",
            "   • Full debate transcript for audit/compliance",
            "   • Human-in-the-loop approval before action",
            "",
            "📊  EXAMPLE",
            "   \"McKinsey Agent recommends pricing response; BCG Agent",
            "    suggests product acceleration; consensus: hybrid approach\""
        ],
        notes="This is where the CaaS marketplace comes alive. The Deep Analysis output feeds into Solution Finding, where multiple branded agents—each trained on a firm's proprietary methodology—debate the options. You get McKinsey's perspective AND BCG's perspective AND your industry expert's perspective, synthesized into a recommendation with full audit trail."
    )
    
    # ========== SLIDE 9: The Pipeline Advantage ==========
    add_table_slide(
        prs,
        "Why the Pipeline Matters",
        ["Capability", "ChatGPT", "Traditional Consulting", "Agent9"],
        [
            ["Proactive monitoring", "❌", "❌", "✅ Always-on"],
            ["Automatic root cause", "❌", "⚠️ Weeks", "✅ Hours"],
            ["Multi-perspective debate", "❌", "❌", "✅ Built-in"],
            ["Evidence-based framing", "❌", "✅", "✅ Automated"],
            ["Audit trail", "❌", "❌", "✅ Complete"],
            ["Human-in-the-loop", "❌", "✅", "✅ Configurable"],
        ],
        notes="Here's why our pipeline is defensible. ChatGPT can't monitor your data—it waits for questions. Traditional consulting is reactive and slow. Only Agent9 combines proactive monitoring, automated investigation, and multi-perspective solution finding with full audit trails."
    )
    
    # ========== SLIDE 10: How It Works (Timeline) ==========
    add_table_slide(
        prs,
        "From Question to Decision in Hours, Not Months",
        ["Step", "Traditional", "Agent9"],
        [
            ["1. Frame Problem", "2-4 weeks", "Minutes"],
            ["2. Gather Data", "4-8 weeks", "Pre-connected"],
            ["3. Analyze", "4-8 weeks", "Hours"],
            ["4. Debate Options", "2-4 weeks", "Real-time"],
            ["5. Recommend", "2-4 weeks", "Instant"],
            ["6. Document", "1-2 weeks", "Automatic"],
            ["TOTAL", "12-24 weeks", "4-24 hours"],
        ],
        notes="Let's break down the timeline. Traditional consulting takes 3-6 months from kickoff to final recommendation. With Agent9, because we're pre-connected to customer data and our agents work in parallel, we compress that to hours."
    )
    
    # ========== SLIDE 6: Platform Capabilities ==========
    add_two_column_slide(
        prs,
        "What We've Built",
        "Core Platform ✅",
        [
            "Multi-agent orchestration",
            "Situation Awareness workflow",
            "Deep Analysis workflow",
            "Solution Finder with LLM debate",
            "Full audit trail & HITL",
            "Data product onboarding"
        ],
        "Coming Soon 🔄",
        [
            "Branded Agent Marketplace",
            "Decision Studio UI",
            "RAG integration for partner IP",
            "Additional data connectors",
            "Partner self-service tools"
        ],
        notes="We're not pitching a vision—we've built the core platform. Our orchestrator coordinates multiple agents through complete workflows. Every decision has an audit trail. What we're building next is the marketplace layer."
    )
    
    # ========== SLIDE 7: Business Model ==========
    add_content_slide(
        prs,
        "Three-Sided Marketplace",
        [
            "PARTNERS (Supply)                    CUSTOMERS (Demand)",
            "• Consulting firms                      • Enterprises",
            "• Tech vendors                           • PE portfolio companies",
            "• Domain experts                        • Mid-market CFOs",
            "",
            "                    ↘    AGENT9 PLATFORM    ↙",
            "",
            "Revenue Streams:",
            "• Platform subscriptions: $80K-$300K/year",
            "• Per-debate fees: $100-$500 each",
            "• Implementation: $50K-$150K",
            "• Partner revenue share: 15-30%"
        ],
        notes="Our business model is a three-sided marketplace. On the supply side, consulting firms and domain experts bring their methodologies. On the demand side, enterprises need strategic insights. Agent9 earns from subscriptions, transaction fees, services, and partner revenue share."
    )
    
    # ========== SLIDE 8: Target Market ==========
    add_table_slide(
        prs,
        "Who We're Targeting",
        ["Segment", "Market Size", "Avg. Deal", "Priority"],
        [
            ["PE Portfolio Ops", "500 firms", "$200K+", "⭐⭐⭐"],
            ["Mid-Market CFOs", "10,000+", "$80K", "⭐⭐⭐"],
            ["Corp Strategy Teams", "500 F500", "$300K+", "⭐⭐"],
            ["Data Platform Customers", "5,000+", "$100K", "⭐⭐"],
        ],
        notes="We're targeting segments that are data-mature and frustrated with traditional consulting. PE portfolio ops teams are ideal—they have a mandate to improve operations and they're tired of paying McKinsey $2M per portfolio company."
    )
    
    # ========== SLIDE 9: Ideal Customer Profile ==========
    add_two_column_slide(
        prs,
        "The Perfect Early Adopter",
        "Must-Haves ✅",
        [
            "Existing data infrastructure",
            "Analytics team to validate",
            "Executive sponsor with budget",
            "$1M+ annual consulting spend"
        ],
        "Trigger Events 🎯",
        [
            "New CFO/CSO hire (first 90 days)",
            "PE acquisition",
            "Activist investor pressure",
            "IPO preparation",
            "Cost reduction initiative"
        ],
        notes="Our ideal customer already has data infrastructure. They have an analytics team who can validate outputs. And they're spending at least a million a year on consulting, so the pain is real."
    )
    
    # ========== SLIDE 10: Go-to-Market ==========
    add_content_slide(
        prs,
        "Go-to-Market: Land, Expand, Partner",
        [
            "PHASE 1: LAND (Q1-Q2 2025)",
            "   • Direct sales to 10-20 early adopters",
            "   • Goal: 5-10 customers, $500K ARR",
            "",
            "PHASE 2: EXPAND (Q3-Q4 2025)",
            "   • Grow within accounts, add use cases",
            "   • Goal: $150K avg ACV, $2M ARR",
            "",
            "PHASE 3: PARTNER (2026)",
            "   • Launch branded agent marketplace",
            "   • Goal: 50% revenue from partner channel"
        ],
        notes="Our go-to-market has three phases. First, we land early adopters through direct sales. Then we expand within those accounts. Finally, we launch the partner marketplace, which becomes our primary growth engine."
    )
    
    # ========== SLIDE 11: Partner Strategy ==========
    add_table_slide(
        prs,
        "Why Partners Will Join",
        ["Partner Fear", "Agent9 Reality"],
        [
            ["\"We'll cannibalize revenue\"", "Access markets you can't serve today"],
            ["\"Our IP will be commoditized\"", "Monetize IP that's currently free"],
            ["\"Clients won't need us\"", "Agents generate leads for human work"],
        ],
        notes="The number one objection from partners is revenue cannibalization. Our reframe: Agent9 lets them access the mid-market—companies that can't afford a $500K engagement. These become qualified leads for human work."
    )
    
    # ========== SLIDE 12: Competitive Landscape ==========
    add_content_slide(
        prs,
        "We're Creating a New Category",
        [
            "                         HIGH EXPERTISE",
            "                               │",
            "        Traditional            │           Agent9",
            "        Consulting             │           CaaS ★",
            "        (McKinsey)             │",
            "                               │",
            "   LOW SCALE ──────────────────┼─────────────── HIGH SCALE",
            "                               │",
            "        Generic AI             │",
            "        (ChatGPT)              │           (empty)",
            "                               │",
            "                         LOW EXPERTISE"
        ],
        notes="Traditional consulting is high expertise but low scale. Generic AI is high scale but low expertise. Agent9 occupies the upper right quadrant: branded expertise at AI scale. No one else is there."
    )
    
    # ========== SLIDE 13: Defensibility ==========
    add_key_points_slide(
        prs,
        "Our Moat Deepens Over Time",
        [
            ("🤝", "Partner Network", "More partners → more expertise → more customers"),
            ("📊", "Data Network", "More customers → better pattern recognition"),
            ("🔒", "Methodology Lock-in", "Partner IP encoded, hard to replicate"),
            ("📋", "Audit Standard", "First to define explainable AI consulting"),
            ("🔄", "Switching Costs", "Data onboarding, workflow integration"),
            ("📈", "Flywheel Effect", "Network effects compound over time"),
        ],
        notes="Our defensibility comes from network effects. As we add partners, we have more expertise, which attracts more customers. Partner methodologies encoded in our platform create lock-in."
    )
    
    # ========== SLIDE 14: Financial Projections ==========
    add_table_slide(
        prs,
        "Path to $10M+ ARR",
        ["Metric", "2025", "2026", "2027"],
        [
            ["Customers", "15", "50", "150"],
            ["Avg. ACV", "$100K", "$120K", "$130K"],
            ["Platform ARR", "$1.5M", "$6M", "$19.5M"],
            ["Partner + Services", "$0.5M", "$1.5M", "$3.5M"],
            ["Total Revenue", "$2M", "$7.5M", "$23M"],
        ],
        notes="We're targeting 15 customers in year one at $100K average deal size, growing to 150 customers by year three. Our unit economics are strong—LTV to CAC ratio of 8-10x with 75-80% gross margins."
    )
    
    # ========== SLIDE 15: Traction ==========
    add_two_column_slide(
        prs,
        "Traction & Milestones",
        "Completed ✅",
        [
            "Multi-agent orchestration",
            "Situation Awareness workflow",
            "Deep Analysis workflow",
            "Solution Finder with LLM debate",
            "Audit trail & HITL checkpoints",
            "Data product onboarding"
        ],
        "Upcoming Milestones",
        [
            "UI demo-ready: Q1 2025",
            "First 3 paying customers: Q1 2025",
            "First partner pilot: Q2 2025",
            "$500K ARR: Q2 2025",
            "Marketplace beta: Q3 2025"
        ],
        notes="We're not just pitching an idea—we have working software. Our core platform is built. We're now polishing the UI for demos and beginning customer pilots."
    )
    
    # ========== SLIDE 16: Team ==========
    add_content_slide(
        prs,
        "The Team",
        [
            "[CEO NAME]",
            "   • [Background in consulting/enterprise software]",
            "",
            "[CTO NAME]",
            "   • [Background in AI/ML, platform architecture]",
            "",
            "[HEAD OF PARTNERSHIPS]",
            "   • [Background in consulting, alliances]",
            "",
            "Advisors:",
            "   • [Former McKinsey Partner]",
            "   • [Former CDO, Fortune 500]",
            "   • [Enterprise SaaS founder]"
        ],
        notes="[Customize with your actual team backgrounds and credentials]"
    )
    
    # ========== SLIDE 17: The Ask ==========
    add_content_slide(
        prs,
        "What We Need to Execute",
        [
            "FOR INVESTORS 💰",
            "   • Raising: $[X]M Seed/Series A",
            "   • Use: 50% Engineering, 30% GTM, 20% Ops",
            "",
            "FOR PARTNERS 🤝",
            "   • Seeking: 2-3 strategic partners",
            "   • Investment: $37.5K + methodology",
            "   • Outcome: Branded agent, revenue share",
            "",
            "FOR CUSTOMERS 🏢",
            "   • Seeking: 5-10 early adopters",
            "   • Investment: $50K-$100K pilot",
            "   • Outcome: On-demand strategic analysis"
        ],
        notes="Here's what we're looking for. From investors, we're raising [X] million. From partners, we want 2-3 strategic firms for co-development. From customers, we're seeking 5-10 early adopter pilots."
    )
    
    # ========== SLIDE 18: Why Now ==========
    add_key_points_slide(
        prs,
        "The Window is Open",
        [
            ("🚀", "Technology Inflection", "LLMs now capable of nuanced reasoning; multi-agent systems production-ready"),
            ("📈", "Market Timing", "Consulting firms under margin pressure; enterprises demand speed"),
            ("🏁", "Competitive Window", "No agentic consulting marketplace exists; first-mover defines category"),
        ],
        notes="Why is now the right time? The technology is ready, the market is ready, and there's no competition. The first mover will define the category."
    )
    
    # ========== SLIDE 19: Summary ==========
    add_key_points_slide(
        prs,
        "Agent9: The Agentic Consulting Marketplace",
        [
            ("💰", "$800B market", "Massive, underserved by technology"),
            ("⚡", "10-50x faster", "Hours, not months"),
            ("📉", "70-90% cheaper", "Accessible to mid-market"),
            ("🗣️", "Multi-perspective", "Best ideas surface through debate"),
            ("✅", "Audit trail", "Compliance-ready, explainable"),
            ("📈", "Platform economics", "Network effects, recurring revenue"),
        ],
        notes="We're going after an $800 billion market with a solution that's 10-50x faster and 70-90% cheaper. We bring multiple perspectives together through AI debate, with full audit trails. We're building the future of consulting."
    )
    
    # ========== SLIDE 20: Thank You ==========
    add_title_slide(
        prs,
        "Thank You",
        "[Your Name] | [email@agent9.ai]",
        "The future of consulting is not human OR AI—it's the best human expertise, encoded and debated by AI, with humans in the loop."
    )
    
    # ========== APPENDIX SECTION ==========
    add_section_slide(prs, "APPENDIX")
    
    # ========== APPENDIX A: Customer ROI ==========
    add_table_slide(
        prs,
        "Appendix A: Customer ROI Detail",
        ["Category", "Before", "After", "Savings"],
        [
            ["Strategy consulting", "$1,000,000", "$200,000", "$800,000"],
            ["Operations consulting", "$900,000", "$150,000", "$750,000"],
            ["Technology advisory", "$800,000", "$100,000", "$700,000"],
            ["Internal analyst support", "$450,000", "$225,000", "$225,000"],
            ["Agent9 platform", "$0", "$180,000", "($180,000)"],
            ["Agent9 usage", "$0", "$120,000", "($120,000)"],
            ["TOTAL", "$3,150,000", "$975,000", "$2,175,000"],
        ],
        notes="Mid-market example with $3M consulting spend. 69% cost reduction with 10x faster insights."
    )
    
    # ========== APPENDIX B: Partner Revenue ==========
    add_table_slide(
        prs,
        "Appendix B: Partner Revenue Model",
        ["Revenue Source", "Partner Share", "Year 1 Example"],
        [
            ["Agent subscription fees", "70-85%", "$350,000"],
            ["Per-debate transaction fees", "60-75%", "$150,000"],
            ["Data onboarding (passive)", "10-15%", "$50,000"],
            ["Human engagement upsells", "100%", "$500,000"],
            ["TOTAL", "", "$1,050,000"],
        ],
        notes="Partners can earn over $1M in year one from a combination of agent fees, transaction revenue, and upsells to human engagements."
    )
    
    # ========== APPENDIX C: Competitive Matrix ==========
    add_table_slide(
        prs,
        "Appendix C: Competitive Feature Matrix",
        ["Feature", "Agent9", "McKinsey", "ChatGPT", "Palantir"],
        [
            ["Strategic expertise", "✅", "✅", "❌", "❌"],
            ["On-demand (24/7)", "✅", "❌", "✅", "✅"],
            ["Multi-agent debate", "✅", "❌", "❌", "❌"],
            ["Full audit trail", "✅", "❌", "❌", "⚠️"],
            ["Enterprise data", "✅", "⚠️", "❌", "✅"],
            ["Cost per insight", "$100-500", "$50K+", "$0.10", "$10K+"],
        ],
        notes="Agent9 is the only solution that combines strategic expertise, on-demand availability, multi-perspective debate, and full audit trails."
    )
    
    return prs


def main():
    """Generate and save the presentation."""
    # Determine output path
    script_dir = Path(__file__).parent
    project_root = script_dir.parent
    output_path = project_root / "docs" / "strategy" / "Agent9_Market_Penetration_Deck_v2.pptx"
    
    # Create presentation
    print("Generating Agent9 Market Penetration Deck...")
    prs = create_presentation()
    
    # Save
    prs.save(str(output_path))
    print(f"✅ Presentation saved to: {output_path}")
    print(f"   Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
